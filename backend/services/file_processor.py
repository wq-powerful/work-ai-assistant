import codecs
import csv
import io
import os
import re
import uuid
from pathlib import Path

MAX_STREAMED_FILE_SIZE = 50 * 1024 * 1024


class FileTooLargeError(ValueError):
    """Raised when an uploaded file exceeds the configured size limit."""


def _build_temp_path(destination: Path) -> Path:
    return destination.parent / f".{destination.name}.{uuid.uuid4().hex}.tmp"


async def stream_upload_to_path(
    upload_file,
    destination: Path,
    *,
    max_bytes: int = MAX_STREAMED_FILE_SIZE,
    chunk_size: int = 1024 * 1024,
) -> int:
    """Stream an uploaded file to disk via a temp file, then atomically replace."""
    destination.parent.mkdir(parents=True, exist_ok=True)
    temp_path = _build_temp_path(destination)
    written = 0

    try:
        with open(temp_path, "wb") as buffer:
            while True:
                chunk = await upload_file.read(chunk_size)
                if not chunk:
                    break

                written += len(chunk)
                if written > max_bytes:
                    raise FileTooLargeError("文件大小超过 50 MiB 限制")

                buffer.write(chunk)

            buffer.flush()
            os.fsync(buffer.fileno())

        os.replace(temp_path, destination)
        return written
    except Exception:
        temp_path.unlink(missing_ok=True)
        raise


def extract_text(file_path: Path) -> str:
    """Extract text content from a file based on its extension."""
    suffix = file_path.suffix.lower()
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".csv": _extract_csv,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xlsx,
        ".pptx": _extract_pptx,
    }
    extractor = extractors.get(suffix)
    if not extractor:
        raise ValueError(f"Unsupported file type: {suffix}")
    raw_text = extractor(file_path)
    return _clean_text(raw_text)


def _clean_text(text: str) -> str:
    """Normalize whitespace and unicode in extracted text."""
    # Normalize unicode
    text = text.replace("\u00a0", " ").replace("\r\n", "\n").replace("\r", "\n")
    # Clean PDF table-of-contents dot leaders (e.g. "Chapter 1 .... 5")
    text = re.sub(r'(?:\.\s*){4,}\d*', ' ', text)
    # Collapse multiple blank lines into two newlines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Collapse multiple spaces into one
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _extract_pdf(file_path: Path) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(str(file_path))
    pages = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            pages.append(page_text)
    return "\n\n".join(pages)


def _extract_docx(file_path: Path) -> str:
    from docx import Document

    doc = Document(str(file_path))
    collected = set()
    paragraphs = []

    def _add_unique(text: str) -> None:
        text = text.strip()
        if text and text not in collected:
            collected.add(text)
            paragraphs.append(text)

    # 1. Paragraphs
    for para in doc.paragraphs:
        _add_unique(para.text)

    # 2. Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            _add_unique(row_text)

    # 3. DrawingML text boxes (<a:txBody> → <a:t>)
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for elem in doc.element.iter():
        if elem.tag.endswith('}txBody') or elem.tag == 'txBody':
            texts = [t.text for t in elem.iter(f'{{{ns_a}}}t') if t.text]
            _add_unique(" ".join(texts))

    # 4. VML text boxes (<v:textbox> → <w:t>)
    ns_w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    ns_v = "urn:schemas-microsoft-com:vml"
    for tb in doc.element.iter(f'{{{ns_v}}}textbox'):
        texts = [t.text for t in tb.iter(f'{{{ns_w}}}t') if t.text]
        _add_unique(" ".join(texts))

    # 5. Headers and footers
    for section in doc.sections:
        for header in [section.header, section.first_page_header, section.even_page_header]:
            if header and header.is_linked_to_previous is False:
                for para in header.paragraphs:
                    _add_unique(para.text)
        for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
            if footer and footer.is_linked_to_previous is False:
                for para in footer.paragraphs:
                    _add_unique(para.text)

    return "\n\n".join(paragraphs)


def _extract_text(file_path: Path) -> str:
    raw_bytes = file_path.read_bytes()
    if not raw_bytes:
        return ""

    bom_encodings = (
        (codecs.BOM_UTF8, "utf-8-sig"),
        (codecs.BOM_UTF16_LE, "utf-16"),
        (codecs.BOM_UTF16_BE, "utf-16"),
    )
    for bom, encoding in bom_encodings:
        if raw_bytes.startswith(bom):
            return raw_bytes.decode(encoding)

    if b"\x00" in raw_bytes:
        even_nulls = raw_bytes[::2].count(0)
        odd_nulls = raw_bytes[1::2].count(0)
        if odd_nulls > even_nulls:
            try:
                return raw_bytes.decode("utf-16-le")
            except UnicodeDecodeError:
                pass
        if even_nulls > odd_nulls:
            try:
                return raw_bytes.decode("utf-16-be")
            except UnicodeDecodeError:
                pass

    for enc in ("utf-8-sig", "gbk", "gb2312"):
        try:
            return raw_bytes.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw_bytes.decode("utf-8", errors="replace")


def _extract_csv(file_path: Path) -> str:
    text_content = _extract_text(file_path)
    reader = csv.reader(io.StringIO(text_content))
    rows = []
    for row in reader:
        rows.append(" | ".join(row))
    return "\n".join(rows)


def _extract_xlsx(file_path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(cell) if cell is not None else "" for cell in row]
            row_text = " | ".join(cell_values)
            if row_text.strip(" |"):
                rows.append(row_text)
        if rows:
            sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets_text)


def _extract_pptx(file_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(paragraph.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)
